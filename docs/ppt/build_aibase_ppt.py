#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""生成 aibase 介绍 PPT（20 页，深色科技风）— 依据 docs/PPT-AIBASE-INTRO-OUTLINE.md

用法: python3 docs/ppt/build_aibase_ppt.py [输出路径]
产物: docs/PPT-AIBASE-INTRO.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ================= 配色 =================
BG     = RGBColor(0x0F, 0x17, 0x2A)   # 深蓝黑
BG2    = RGBColor(0x1B, 0x24, 0x3B)   # 次级背景
CARD   = RGBColor(0x16, 0x20, 0x38)   # 卡片
WHITE  = RGBColor(0xF8, 0xFA, 0xFC)
MUTED  = RGBColor(0x94, 0xA3, 0xB8)
INDIGO = RGBColor(0x63, 0x66, 0xF1)
CYAN   = RGBColor(0x22, 0xD3, 0xEE)
GOLD   = RGBColor(0xF5, 0x9E, 0x0B)
RED    = RGBColor(0xF8, 0x71, 0x71)
GREEN  = RGBColor(0x34, 0xD3, 0x99)
GRID   = RGBColor(0x2A, 0x35, 0x52)
FONT   = "Microsoft YaHei"
MONO   = "Consolas"

SW, SH = Inches(13.333), Inches(7.5)
TOTAL = 20

# ================= 基础工具 =================

def new_prs():
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH
    return prs

def add_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def set_bg(slide, color=BG):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color

def add_rect(slide, x, y, w, h, fill=None, line=None, line_w=None, shape=MSO_SHAPE.RECTANGLE):
    sp = slide.shapes.add_shape(shape, x, y, w, h)
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = line_w or Pt(1)
    sp.shadow.inherit = False
    return sp

def add_text(slide, x, y, w, h, text, size=18, color=WHITE, bold=False,
             align=PP_ALIGN.LEFT, font=FONT, anchor=MSO_ANCHOR.TOP, wrap=True,
             spacing=1.0, space_after=0, italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    for i, ln in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        p.space_after = Pt(space_after)
        r = p.add_run(); r.text = ln
        f = r.font
        f.size = Pt(size); f.bold = bold; f.italic = italic
        f.color.rgb = color; f.name = font
    return tb

def add_rich(slide, x, y, w, h, items, anchor=MSO_ANCHOR.TOP, spacing=1.12, space_after=7):
    """items: list of (text, size, color, bold, bullet)"""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    for i, (text, size, color, bold, bullet) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = spacing; p.space_after = Pt(space_after)
        if bullet:
            rb = p.add_run(); rb.text = "▪ "
            rb.font.size = Pt(size); rb.font.bold = True
            rb.font.color.rgb = GOLD; rb.font.name = FONT
        r = p.add_run(); r.text = text
        f = r.font
        f.size = Pt(size); f.bold = bold; f.color.rgb = color; f.name = FONT
    return tb

def add_card(slide, x, y, w, h, title=None, body=None, accent=INDIGO,
             title_size=16, body_size=13, body_color=MUTED):
    add_rect(slide, x, y, w, h, fill=CARD, line=GRID, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    if accent is not None:
        add_rect(slide, x, y + Inches(0.14), Pt(4), h - Inches(0.28), fill=accent)
    if title:
        add_text(slide, x + Inches(0.22), y + Inches(0.12), w - Inches(0.4), Inches(0.42),
                 title, size=title_size, bold=True, color=WHITE)
    if body:
        ty = y + Inches(0.58) if title else y + Inches(0.16)
        add_text(slide, x + Inches(0.22), ty, w - Inches(0.4), h - Inches(0.72),
                 body, size=body_size, color=body_color, spacing=1.18)

def add_badge(slide, x, y, w, h, text, fill, size=13, bold=True, color=WHITE):
    add_rect(slide, x, y, w, h, fill=fill, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text(slide, x, y + Inches(0.04), w, h - Inches(0.06), text,
             size=size, bold=bold, color=color, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

def add_table(slide, x, y, w, headers, rows, col_w=None, header_fill=INDIGO,
              row_h=Inches(0.44), header_h=Inches(0.5), font_size=13, header_size=14):
    n_rows, n_cols = len(rows) + 1, len(headers)
    gt = slide.shapes.add_table(n_rows, n_cols, x, y, w, Inches(0.5) * n_rows).table
    if col_w:
        for i, cw in enumerate(col_w):
            gt.columns[i].width = cw
    gt.rows[0].height = header_h
    for i in range(1, n_rows):
        gt.rows[i].height = row_h
    for j, htext in enumerate(headers):
        cell = gt.cell(0, j)
        cell.fill.solid(); cell.fill.fore_color.rgb = header_fill
        cell.margin_top = Pt(3); cell.margin_bottom = Pt(3)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf = cell.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
        r = p.add_run(); r.text = htext
        r.font.size = Pt(header_size); r.font.bold = True
        r.font.color.rgb = WHITE; r.font.name = FONT
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            cell = gt.cell(i, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = CARD if i % 2 == 1 else BG2
            cell.margin_top = Pt(3); cell.margin_bottom = Pt(3)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf = cell.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
            r = p.add_run(); r.text = str(val)
            r.font.size = Pt(font_size); r.font.color.rgb = WHITE; r.font.name = FONT
    return gt

def header(slide, title, subtitle=None, n=None, accent=GOLD):
    add_rect(slide, Inches(0.6), Inches(0.5), Inches(0.12), Inches(0.62), fill=accent)
    add_text(slide, Inches(0.9), Inches(0.44), Inches(11.4), Inches(0.7),
             title, size=28, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, Inches(0.92), Inches(1.08), Inches(11.6), Inches(0.4),
                 subtitle, size=13.5, color=MUTED)
    add_rect(slide, Inches(0.6), SH - Inches(0.6), SW - Inches(1.2), Pt(1), fill=GRID)
    # 右上角装饰
    add_rect(slide, SW - Inches(1.0), Inches(0.45), Inches(0.28), Inches(0.28),
             fill=accent, shape=MSO_SHAPE.OVAL)
    add_rect(slide, SW - Inches(0.62), Inches(0.63), Inches(0.16), Inches(0.16),
             fill=CYAN, shape=MSO_SHAPE.OVAL)
    if n:
        add_text(slide, SW - Inches(1.4), SH - Inches(0.52), Inches(1.0), Inches(0.3),
                 f"{n:02d} / {TOTAL}", size=11, color=MUTED, align=PP_ALIGN.RIGHT)

def arrow(slide, x, y, w, h, color=CYAN, direction="right"):
    shp = MSO_SHAPE.RIGHT_ARROW if direction == "right" else MSO_SHAPE.DOWN_ARROW
    add_rect(slide, x, y, w, h, fill=color, shape=shp)

# ================= 页面 =================

def slide_01(prs):
    s = add_slide(prs); set_bg(s)
    # 背景装饰：右上大圆、左下小圆、顶部光带
    add_rect(s, SW - Inches(2.6), -Inches(1.4), Inches(4.2), Inches(4.2), fill=INDIGO, shape=MSO_SHAPE.OVAL)
    add_rect(s, SW - Inches(1.9), -Inches(0.9), Inches(3.0), Inches(3.0), fill=BG, shape=MSO_SHAPE.OVAL)
    add_rect(s, -Inches(1.0), SH - Inches(2.2), Inches(3.2), Inches(3.2), fill=CYAN, shape=MSO_SHAPE.OVAL)
    add_rect(s, Inches(0.6), Inches(0.55), Inches(2.4), Pt(3), fill=GOLD)
    # 主标题
    add_text(s, Inches(0.9), Inches(1.7), Inches(9), Inches(1.4), "aibase",
             size=88, bold=True, color=WHITE)
    add_text(s, Inches(0.95), Inches(3.1), Inches(11.5), Inches(0.6),
             "AI Agent Operating System + Reliable AI Engineering Framework",
             size=22, bold=True, color=CYAN)
    add_text(s, Inches(0.95), Inches(3.85), Inches(11.5), Inches(0.55),
             "让 AI 从「会写代码」到「可靠地做成事」", size=18, color=MUTED)
    # 五大创新标签
    tags = ["AIOS 内核", "角色式 Agent", "治理协议", "可靠执行闭环", "评估体系"]
    x = Inches(0.95)
    for i, t in enumerate(tags):
        w = Inches(1.72)
        col = [INDIGO, CYAN, GOLD, GREEN, RED][i]
        add_badge(s, x, Inches(4.7), w, Inches(0.5), t, col, size=14)
        x += w + Inches(0.15)
    # 底部署名
    add_rect(s, Inches(0.95), Inches(5.75), Inches(0.6), Pt(2), fill=GRID)
    add_text(s, Inches(0.95), Inches(6.0), Inches(8), Inches(0.4),
             "hb  ·  MIT License  ·  2026", size=14, color=MUTED)
    add_text(s, Inches(0.95), Inches(6.45), Inches(8), Inches(0.4),
             "「让 AI 干活，先给 AI 立规矩」", size=13, color=GRID, italic=True)

def slide_02(prs):
    s = add_slide(prs); set_bg(s)
    header(s, "为什么需要它？", "第一代 AI 模板假设「AI = 会读规则的助手」，2026 年的 AI agent 需要的是「操作系统」。", 2)
    # 左：旧假设
    add_card(s, Inches(0.6), Inches(1.7), Inches(5.9), Inches(2.5),
             "第一代假设", "AI 工具 → 读规则 → 执行 → 完成\n\n「AI 本质上是会读规则的助手」", accent=GREEN)
    # 右：新现实
    add_card(s, Inches(6.8), Inches(1.7), Inches(5.9), Inches(2.5),
             "2026 现实", "AI 已能自主规划、多步执行、调用工具\n但没有治理约束时——", accent=RED)
    # 四个失控点
    pts = [("乱改不该改的文件", "生成目录"), ("忽略系统架构", "不在上下文里"),
           ("无法自评", "AI 不能审查自己"), ("失败无修复", "验证挂了没人管")]
    xs = [Inches(0.6), Inches(3.4), Inches(6.8), Inches(9.6)]
    for (a, b), x in zip(pts, xs):
        add_card(s, x, Inches(4.5), Inches(3.1), Inches(1.05),
                 "✗ " + a, b, accent=RED, title_size=13.5, body_size=11)
    # 结论
    add_rect(s, Inches(0.6), Inches(5.9), Inches(12.1), Inches(1.0), fill=BG2, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text(s, Inches(0.9), Inches(6.05), Inches(11.6), Inches(0.7),
             "结论：需要「AI 操作系统」—— 理解目标 → 建模 → 规划 → 执行 → 验证 → 学习",
             size=16, bold=True, color=GOLD, anchor=MSO_ANCHOR.MIDDLE)

def slide_03(prs):
    s = add_slide(prs); set_bg(s)
    header(s, "aibase 是什么？", "不是「AI 配置模板」，是 AI Agent Operating System —— 完整工程治理框架。", 3)
    add_table(s, Inches(0.6), Inches(1.7), Inches(12.1),
              ["#", "核心创新", "说明"],
              [["1", "🧠 AIOS 内核", "Governance + Cognition + Execution + Context + Memory"],
               ["2", "🤖 角色式 Agent", "按角色（Manager/Coder/Reviewer…）而非按工具定义"],
               ["3", "🔒 Governance 协议", "明确「AI 不可以做什么」的边界"],
               ["4", "🔁 可靠执行闭环", "Plan → Impact → Execute → Reflect → Verify → Repair"],
               ["5", "📊 评估体系", "任务成功率 / 返工 / Token 消耗 / 模型对比"]],
              col_w=[Inches(0.8), Inches(3.6), Inches(7.7)], row_h=Inches(0.78))
    add_text(s, Inches(0.6), Inches(6.35), Inches(12), Inches(0.5),
             "一句话：让 AI 从「收到命令就执行」变成「理解目标、可靠交付」的执行者。",
             size=15, bold=True, color=CYAN)

def slide_04(prs):
    s = add_slide(prs); set_bg(s)
    header(s, "整体架构 · kit 布局", "aibase 自身就是 kit 布局（自举 / dogfood），升级 = 整体替换 kit/ 目录。", 4)
    lines = [
        ("AGENTS.md", "项目入口 —— 25+ AI 工具原生读取"),
        ("kit/aios/", "AIOS 内核：governance / cognition / execution / context / memory"),
        ("kit/agents/", "7 个角色定义（Manager / Coder / Reviewer…）"),
        ("kit/profiles/", "项目类型模板：backend / game-server / unity / unreal / frontend / data…"),
        ("kit/knowledge/", "项目知识图谱（架构 / 模块 / 依赖 / ADR / 术语）"),
        ("kit/cli/", "统一控制入口：task / mkproject / persona / protect / sandbox"),
        ("kit/evaluation/", "质量评估体系（metrics / failures / reports / benchmarks）"),
        ("runtime/", "文件即数据库 —— 任务 / 状态 / 验证 / 审查全在文件系统"),
    ]
    y = Inches(1.6)
    for name, desc in lines:
        add_rect(s, Inches(0.6), y, Inches(3.3), Inches(0.52), fill=BG2, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        add_text(s, Inches(0.78), y + Inches(0.06), Inches(3.1), Inches(0.4),
                 name, size=15, bold=True, color=CYAN)
        add_text(s, Inches(4.15), y + Inches(0.06), Inches(8.6), Inches(0.4),
                 desc, size=14, color=WHITE)
        y += Inches(0.59)
    add_badge(s, Inches(9.4), Inches(6.5), Inches(3.3), Inches(0.5), "零依赖 · 可移植 · 文件即数据库", GOLD, size=13)

def slide_05(prs):
    s = add_slide(prs); set_bg(s)
    header(s, "治理层 · AI 不可以做什么", "治理先于执行 —— 任何 AI 动作前先过 governance 检查。", 5)
    add_table(s, Inches(0.6), Inches(1.7), Inches(12.1),
              ["协议", "管什么"],
              [["task-policy", "任务状态机 open→in-progress→in-review→done；分级治理 fast-path；返工上限"],
               ["modification-policy", "目录可写/禁写（generated_dirs 禁写）；高风险文件分级"],
               ["security-policy", "Rule of Two；密钥不进提示词；训练数据退出"],
               ["risk-policy", "风险分级 P0-P3 与审批条件"],
               ["roadmap-policy", "路线图唯一权威与变更控制"]],
              col_w=[Inches(3.4), Inches(8.7)], row_h=Inches(0.68))
    add_text(s, Inches(0.6), Inches(6.1), Inches(12), Inches(0.6),
             "金句：AI 最需要知道的不是「能做什么」，而是「不能做什么」。",
             size=16, bold=True, color=GOLD)

def slide_06(prs):
    s = add_slide(prs); set_bg(s)
    header(s, "风险分级与权限模型", "治理强度随风险分级 —— 不搞一刀切。", 6)
    # 四级风险卡片
    lv = [("🔴 P0", "不可逆 / 高影响", "数据库 schema · 支付 · 账号 · 权限", "必须人工批准"),
          ("🟡 P1", "可逆但需审查", "API 契约 · 共享类型 · CI/CD", "TASK 声明 + reviewer 审查"),
          ("🟢 P2", "常规开发", "业务逻辑 · UI · 测试 · 文档", "TASK 追踪"),
          ("⚪ P3", "低风险", "注释 · 格式 · 命名 · README", "自由执行")]
    col = [RED, GOLD, GREEN, MUTED]
    xs = [Inches(0.6), Inches(3.8), Inches(7.0), Inches(10.2)]
    for (tag, name, ex, rule), x, c in zip(lv, xs, col):
        add_card(s, x, Inches(1.7), Inches(3.05), Inches(2.5),
                 tag + "  " + name, ex + "\n\n" + rule, accent=c,
                 title_size=15, body_size=11.5)
    # 自动升级
    add_card(s, Inches(0.6), Inches(4.5), Inches(12.1), Inches(1.05),
             "风险自动升级（任一条件触发 +1 级）",
             "改动 >5 文件   ·   跨 >2 模块   ·   涉及 >1 种语言   ·   存在未解决依赖",
             accent=INDIGO, title_size=14, body_size=13)
    # 修改权限
    add_rect(s, Inches(0.6), Inches(5.85), Inches(12.1), Inches(0.95), fill=BG2, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text(s, Inches(0.9), Inches(5.95), Inches(11.6), Inches(0.8),
             "修改权限：generated_dirs 禁写（改生成代码 = 改规格 → 重新生成）；只有 source_dirs 可写",
             size=15, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)

def slide_07(prs):
    s = add_slide(prs); set_bg(s)
    header(s, "执行闭环 · Reliable Execution Loop", "验证是流程的一部分，不是 AI 的自觉。", 7, accent=CYAN)
    # §0 分流
    add_badge(s, Inches(0.6), Inches(1.55), Inches(3.1), Inches(0.5), "§0 启动分流", INDIGO, size=14)
    add_text(s, Inches(3.9), Inches(1.6), Inches(8.8), Inches(0.5),
             "改契约 / 规格 / 生成源 / 跨端协议？ → 是：SDD 分支（规格先行）｜ 否：通用路径",
             size=13, color=MUTED)
    # 六环
    steps = [("§1 Plan", "计划先行", INDIGO), ("§2 Impact", "影响分析", INDIGO),
             ("§3 Execute", "按计划执行", INDIGO), ("§4 Reflect", "自省清单", INDIGO),
             ("§5 Verify", "机械强制", GREEN), ("§6 Repair", "返工修复", GOLD)]
    x = Inches(0.6)
    for name, desc, c in steps:
        add_card(s, x, Inches(2.4), Inches(1.86), Inches(1.55),
                 name, desc, accent=c, title_size=15, body_size=12)
        if name != "§6 Repair":
            arrow(s, x + Inches(1.88), Inches(2.95), Inches(0.3), Inches(0.42), color=GRID)
        x += Inches(2.06)
    # 说明
    add_rich(s, Inches(0.6), Inches(4.3), Inches(12.1), Inches(2.3),
             [("§5 Verify 机械强制：", 14, WHITE, True, True),
              ("build / lint / test / check 真跑，cli/task verify 通过才写 VERIFY 记录（不是自证）", 14, MUTED, False, False),
              ("§6 Repair：", 14, WHITE, True, True),
              ("失败 → 根因分析 → 修复；返工上限 3 次，超限人工介入 —— 不无限烧循环", 14, MUTED, False, False),
              ("fast-path：", 14, WHITE, True, True),
              ("低风险任务可跳过部分仪式，但 verify 与 TASK 状态流转对任何任务不可省略", 14, MUTED, False, False)],
             spacing=1.15, space_after=8)

def slide_08(prs):
    s = add_slide(prs); set_bg(s)
    header(s, "SDD 规格驱动开发", "规格是唯一真相 —— 规格变更先于实现。", 8, accent=INDIGO)
    add_rich(s, Inches(0.6), Inches(1.7), Inches(5.8), Inches(3.4),
             [("三种项目形态", 16, CYAN, True, True),
              ("contract —— API 契约 + 代码生成", 14, WHITE, False, False),
              ("docs —— 工具 / 游戏客户端（文档即规格）", 14, WHITE, False, False),
              ("protocol —— CS 网游（客户端+服务器共同契约）", 14, WHITE, False, False),
              ("", 6, MUTED, False, False),
              ("五个场景", 16, CYAN, True, True),
              ("A / A' / A\" —— 规格变更（先改规格→校验→生成→实现）", 14, WHITE, False, False),
              ("B —— 仅逻辑修改", 14, WHITE, False, False),
              ("C —— Bug 修复", 14, WHITE, False, False)],
             spacing=1.1, space_after=6)
    add_card(s, Inches(6.9), Inches(1.7), Inches(5.8), Inches(3.3),
             "核心约束", "规格是唯一真相\n生成代码不可手动编辑\n规格变更先于实现\n验证是流程的一部分\n任务可追踪", accent=GOLD)
    add_rect(s, Inches(0.6), Inches(5.4), Inches(12.1), Inches(1.1), fill=BG2, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text(s, Inches(0.9), Inches(5.55), Inches(11.6), Inches(0.8),
             "判断不清时 → 先按规格影响处理，并在 TASK 中记录假设。",
             size=15, bold=True, color=CYAN, anchor=MSO_ANCHOR.MIDDLE)

def slide_09(prs):
    s = add_slide(prs); set_bg(s)
    header(s, "角色体系 · 按角色而非工具", "新范式：Coder 角色需要什么能力？ —— 而不是「Claude 能做什么？」", 9)
    add_table(s, Inches(0.6), Inches(1.7), Inches(12.1),
              ["角色", "职责", "写代码", "关键约束"],
              [["Manager", "任务分派、进度跟踪、阻塞管理", "✗", "只管流程"],
               ["Analyst", "需求分析、影响评估、假设管理", "✗", "不明不写"],
               ["Architect", "架构决策、模块划分、技术选型", "✗", "ADR 记录"],
               ["Coder", "按计划实现代码", "✓", "必须通过 verify"],
               ["Reviewer", "代码审查、安全检查", "✗", "生成者 ≠ 审查者"],
               ["Tester", "测试编写与执行", "✓(测试)", "测试先行"],
               ["Researcher", "外部检索、知识录入", "✗", "写入 knowledge/"]],
              col_w=[Inches(2.0), Inches(5.1), Inches(1.2), Inches(3.8)], row_h=Inches(0.55))
    add_text(s, Inches(0.6), Inches(6.35), Inches(12), Inches(0.5),
             "新 AI 工具（Codex / Gemini CLI / Qwen Code…）不断出现 —— 按角色定义，无需为每个工具写适配器。",
             size=15, bold=True, color=CYAN)

def slide_10(prs):
    s = add_slide(prs); set_bg(s)
    header(s, "机械强制层 · 四道墙", "规则写在文档里是建议，写在文件系统 / Git / CI 里才是纪律。", 10)
    add_table(s, Inches(0.6), Inches(1.7), Inches(12.1),
              ["层", "机制", "挡什么"],
              [["进程沙箱", "sandbox-run（无网络容器）", "Rule of Two：敏感数据不可外传"],
               ["文件权限", "protect（chmod 锁 generated_dirs 只读）", "手动改生成代码"],
               ["Git hook", "pre-commit / commit-msg", "无 TASK 引用的提交 · 碰 generated_dirs"],
               ["CI", "GitHub Actions verify.yml", "绕过本地 hook 也挡得住"]],
              col_w=[Inches(2.4), Inches(5.2), Inches(4.5)], row_h=Inches(0.7))
    add_text(s, Inches(0.6), Inches(5.6), Inches(12), Inches(0.8),
             "这四道墙跟用 Claude Code、Cursor、pi 还是人手改代码无关 —— 都会经过文件系统、git、CI。\n这是 aibase 与一般 prompt 模板的本质区别：不靠 AI 自觉，靠系统强制。",
             size=15, bold=True, color=GOLD, spacing=1.3)

def slide_11(prs):
    s = add_slide(prs); set_bg(s)
    header(s, "CLI 工具箱", "统一控制入口：已自动化的流程有命令；未自动化的流程显式读写 Markdown 记录。", 11)
    cmds = [("mkproject", "用 kit 布局创建新项目（--profile / --persona）", CYAN),
            ("task new / list / start", "任务生命周期管理", CYAN),
            ("task verify", "真实执行 build / lint / test / check，通过才写 VERIFY", GREEN),
            ("task review / approve / done", "审查与关闭（fast-path 可跳 review）", INDIGO),
            ("persona list / use / off", "人格系统（表达层设定）", GOLD),
            ("check / protect / sandbox-run", "健康检查 / 锁生成目录 / 无网络沙箱", INDIGO),
            ("autoloop-coder / reviewer", "无人值守循环（实验性，隔离环境）", GOLD)]
    y = Inches(1.65)
    for name, desc, c in cmds:
        add_badge(s, Inches(0.6), y, Inches(3.4), Inches(0.5), name, c, size=13)
        add_text(s, Inches(4.3), y + Inches(0.05), Inches(8.4), Inches(0.45), desc, size=14, color=WHITE)
        y += Inches(0.62)
    add_text(s, Inches(0.6), Inches(6.2), Inches(12), Inches(0.5),
             "⚠️ --unattended 会给 claude 传 --dangerously-skip-permissions —— 仅限隔离环境（容器 / worktree / 沙箱）。",
             size=13, color=RED)

def slide_12(prs):
    s = add_slide(prs); set_bg(s)
    header(s, "知识图谱 · knowledge/", "AI 最大的问题不是不会写代码，而是不知道系统关系。", 12)
    kp = [("architecture/", "系统架构设计（分层/模块/通信）"),
          ("modules/", "模块清单：职责、依赖、入口"),
          ("dependencies/", "模块间依赖关系图"),
          ("decisions/", "ADR —— 为什么这样设计"),
          ("history/", "已知问题、技术债、失败尝试"),
          ("glossary/", "领域术语表（统一语言）")]
    xs = [Inches(0.6), Inches(4.65), Inches(8.7)]
    ys = [Inches(1.7), Inches(3.7)]
    for i, (name, desc) in enumerate(kp):
        x = xs[i % 3]; y = ys[i // 3]
        add_card(s, x, y, Inches(3.85), Inches(1.85), name, desc,
                 accent=[INDIGO, CYAN, GOLD][i % 3], title_size=15, body_size=12.5)
    add_rect(s, Inches(0.6), Inches(5.85), Inches(12.1), Inches(1.0), fill=BG2, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text(s, Inches(0.9), Inches(5.95), Inches(11.6), Inches(0.8),
             "原则：Knowledge invisible to the agent doesn't exist.  ·  Module before change（改代码前先查模块文档）",
             size=15, bold=True, color=CYAN, anchor=MSO_ANCHOR.MIDDLE)

def slide_13(prs):
    s = add_slide(prs); set_bg(s)
    header(s, "评估体系 · evaluation/", "不评估 = 不知道 AI 做得好不好。", 13)
    add_table(s, Inches(0.6), Inches(1.7), Inches(12.1),
              ["维度", "指标", "记录位置"],
              [["任务成功率", "完成数 / 分配数；done 率；cancelled 率", "metrics/"],
               ["返工次数", "in-review → in-progress 打回次数 / 任务", "failures/"],
               ["Bug 率", "审查发现问题数 / 变更文件数", "metrics/"],
               ["Token 消耗", "平均 token / 任务；超额 token", "metrics/"],
               ["耗时", "平均 wall time / 任务", "reports/"],
               ["模型表现", "按模型分类的成功率 / 失败原因", "benchmarks/"]],
              col_w=[Inches(2.6), Inches(7.0), Inches(2.5)], row_h=Inches(0.6))
    add_text(s, Inches(0.6), Inches(6.0), Inches(12), Inches(0.6),
             "示例：同一任务 Claude Sonnet vs GPT-4o —— 横向对比模型能力。",
             size=15, bold=True, color=GOLD)

def slide_14(prs):
    s = add_slide(prs); set_bg(s)
    header(s, "人格系统 · 表达层设定", "人格只改语气，不改结论 —— 治理纪律不变。", 14, accent=GOLD)
    add_rich(s, Inches(0.6), Inches(1.7), Inches(6.6), Inches(4.0),
             [("表达层 ≠ 认知层", 16, CYAN, True, True),
              ("语气可以疯、可以妖、可以狂，结论必须准", 14, WHITE, False, False),
              ("", 6, MUTED, False, False),
              ("按需加载", 16, CYAN, True, True),
              ("personas/active.md 存在则激活，缺失则零加载", 14, WHITE, False, False),
              ("", 6, MUTED, False, False),
              ("单源真相 + 薄壳指路", 16, CYAN, True, True),
              ("人格内容唯一真相在人格库，各工具薄壳只指路不复制", 14, WHITE, False, False),
              ("", 6, MUTED, False, False),
              ("严谨模式", 16, CYAN, True, True),
              ("随时可去人格化，回归理性直白表达", 14, WHITE, False, False)],
             spacing=1.1, space_after=6)
    add_card(s, Inches(7.6), Inches(1.7), Inches(5.1), Inches(4.0),
             "人格库（部分）", "东方不败 · 曹操 · 孙悟空 · 林黛玉\n福尔摩斯 · 蝙蝠侠 · 东方不败 · 哪吒\n…… 50+ 人格待命",
             accent=GOLD, title_size=16, body_size=14)

def slide_15(prs):
    s = add_slide(prs); set_bg(s)
    header(s, "快速上手 · 5 分钟 Demo", "从零到第一个任务，只需 8 条命令。", 15, accent=GREEN)
    code = ("git clone <aibase> /tmp/framework\n"
            "python /tmp/framework/kit/cli/mkproject ~/code/my-project --profile backend\n"
            "bash kit/cli/protect\n"
            "python kit/cli/task new \"实现用户登录\" --priority P1 --reviewer claude\n"
            "python kit/cli/task start TASK-001\n"
            "python kit/cli/task verify TASK-001   # 真跑 build/lint/test/check\n"
            "python kit/cli/task review TASK-001\n"
            "python kit/cli/task approve TASK-001")
    add_rect(s, Inches(0.6), Inches(1.7), Inches(12.1), Inches(4.4), fill=BG2, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_rect(s, Inches(0.6), Inches(1.7), Inches(12.1), Inches(0.45), fill=INDIGO, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text(s, Inches(0.9), Inches(1.76), Inches(6), Inches(0.35), "Terminal — aibase quickstart", size=12, bold=True, color=WHITE)
    add_text(s, Inches(0.9), Inches(2.35), Inches(11.5), Inches(3.6), code,
             size=14.5, color=CYAN, font=MONO, spacing=1.45)
    add_text(s, Inches(0.6), Inches(6.35), Inches(12), Inches(0.5),
             "重点：task verify 是「真跑命令」而不是手写记录 —— 通过才生成 VERIFY 记录。",
             size=15, bold=True, color=GOLD)

def slide_16(prs):
    s = add_slide(prs); set_bg(s)
    header(s, "实际应用 · dogfood 实证", "框架自己就是自己的第一个用户。", 16, accent=GREEN)
    add_card(s, Inches(0.6), Inches(1.7), Inches(5.9), Inches(2.2),
             "aibase 自举", "自身即 kit 布局\n58 个 TASK —— 任务 / 审查 / 验证记录齐全\n升级 = 整体替换 kit/ 目录", accent=GREEN, title_size=16, body_size=13.5)
    add_card(s, Inches(6.8), Inches(1.7), Inches(5.9), Inches(2.2),
             "已同步注册项目", "aimonitor · account · baseline\nx1-prototype · x1design · westhill\n（6+ 个 AIOS 项目）", accent=INDIGO, title_size=16, body_size=13.5)
    add_rich(s, Inches(0.6), Inches(4.2), Inches(12.1), Inches(2.0),
             [("真实运行证据：", 15, WHITE, True, True),
              ("cli/task list 状态机流转  ·  REVIEW 审查记录  ·  VERIFY 验证记录  ·  autoloop 心跳与事件流", 14, MUTED, False, False),
              ("aimonitor 正在实时监控这些项目 —— 下一步揭晓", 14, CYAN, False, False)],
             spacing=1.2, space_after=8)

def slide_17(prs):
    s = add_slide(prs); set_bg(s)
    header(s, "aimonitor · 框架造出的第一个成品", "实时监控多个 AI 项目的仪表盘，也是框架设计的「活体验证」。", 17, accent=CYAN)
    # 架构图三栏
    add_card(s, Inches(0.6), Inches(1.7), Inches(3.6), Inches(2.0),
             "被监控项目", "7+ 个 AIOS 项目\nruntime/ 文件\n任务 · 心跳 · 事件 · VERIFY", accent=INDIGO, title_size=14, body_size=11.5)
    arrow(s, Inches(4.32), Inches(2.55), Inches(0.42), Inches(0.42), color=GOLD)
    add_card(s, Inches(4.85), Inches(1.7), Inches(3.6), Inches(2.0),
             "采集层", "local：直读文件系统\nagent：遥测组件随\nmkproject 自动分发", accent=GOLD, title_size=14, body_size=11.5)
    arrow(s, Inches(8.57), Inches(2.55), Inches(0.42), Inches(0.42), color=GOLD)
    add_card(s, Inches(9.1), Inches(1.7), Inches(3.6), Inches(2.0),
             "aimonitor 仪表盘", "网页实时展示\n任务全景 · 心跳存活\n事件流 · 验证审查", accent=CYAN, title_size=14, body_size=11.5)
    # 四张牌
    pts = [("① 能造真成品", "前后端 + 多机部署 + 运维手册 + 注册审批"),
           ("② 文件即数据库被消费", "读的正是 runtime/ 里的全部产物"),
           ("③ 跨项目已运行", "监控 7+ 项目，含远程 Windows/Linux"),
           ("④ 零依赖可行", "Python 标准库 + 手写前端")]
    xs = [Inches(0.6), Inches(3.7), Inches(6.8), Inches(9.9)]
    for (a, b), x in zip(pts, xs):
        add_card(s, x, Inches(4.0), Inches(2.95), Inches(1.5), a, b,
                 accent=[GREEN, INDIGO, CYAN, GOLD][pts.index((a, b))], title_size=13.5, body_size=10.5)
    add_rect(s, Inches(0.6), Inches(5.8), Inches(12.1), Inches(1.0), fill=BG2, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text(s, Inches(0.9), Inches(5.9), Inches(11.6), Inches(0.8),
             "彩蛋：它自己也是用这套框架开发的 —— 69 个 TASK + MONITOR-SPEC 规格驱动。既是框架的产品，也是框架的镜子。",
             size=14, bold=True, color=GOLD, anchor=MSO_ANCHOR.MIDDLE)

def slide_18(prs):
    s = add_slide(prs); set_bg(s)
    header(s, "演进路线 · 从单项目自动到跨项目流水线", "AI 负责跑，人负责定方向、批风险、验收果。", 18, accent=INDIGO)
    add_table(s, Inches(0.6), Inches(1.7), Inches(12.1),
              ["阶段", "内容", "现状"],
              [["① 单项目自动闭环", "autoloop coder+reviewer；P0/返工超限转人工", "✅ 已实现（实验性）"],
               ["② 多项目监控", "agent 遥测 + aimonitor 仪表盘；心跳/告警", "✅ 已实现"],
               ["③ 多项目自动调度", "调度器读多项目 runtime/tasks/ 分配 coder/reviewer", "🟡 架构就绪，缺编排层"],
               ["④ 点子→成品流水线", "跨项目 DAG + 阶段人工闸门 + 成品验收", "🔵 愿景，需逐步验证"]],
              col_w=[Inches(3.2), Inches(6.6), Inches(2.3)], row_h=Inches(0.72))
    add_rich(s, Inches(0.6), Inches(5.15), Inches(12.1), Inches(1.6),
             [("已内建的人机分工点：", 14, WHITE, True, True),
              ("P0 风险操作 → 必须人工批准（否则 blocked）  ·  返工 ≥3 次 → 机械拒绝  ·  卡死/阻塞 → task block + 告警  ·  成品验收仍需人",
               13, MUTED, False, False)],
             spacing=1.2, space_after=6)

def slide_19(prs):
    s = add_slide(prs); set_bg(s)
    header(s, "与生态标准的关系", "不是又一个孤立框架，是与标准生态对接的工程层。", 19, accent=GREEN)
    rows = [
        ["AGENTS.md", "标准格式，25+ AI 工具原生接入，零配置（Claude Code 用 symlink）"],
        ["AIOS (agiresearch)", "设计理念同源（Governance/Context/Memory/Storage），作为文件级治理层独立实现"],
        ["MCP", "kit/tools/ 能力层可包装为 MCP server，agent 通过标准协议调用"],
        ["A2A", "kit/agents/ 角色可通过 A2A 协议跨框架协作（远期）"],
        ["跨平台", "Windows（PowerShell + python）/ Linux / macOS"],
    ]
    add_table(s, Inches(0.6), Inches(1.7), Inches(12.1),
              ["标准", "关系"], rows,
              col_w=[Inches(3.4), Inches(8.7)], row_h=Inches(0.82))
    add_text(s, Inches(0.6), Inches(6.3), Inches(12), Inches(0.5),
             "所有命令零网络依赖（autoloop-* 例外）—— 框架本身不吃云、不锁厂商。",
             size=15, bold=True, color=CYAN)

def slide_20(prs):
    s = add_slide(prs); set_bg(s)
    header(s, "总结 · 可靠 AI 工程", None, 20, accent=GOLD)
    add_text(s, Inches(0.6), Inches(1.6), Inches(12.1), Inches(0.8),
             "治理 + 角色 + 闭环 + 机械强制 + 知识 + 评估 = 可靠 AI 工程",
             size=26, bold=True, color=GOLD)
    add_rich(s, Inches(0.6), Inches(2.7), Inches(12.1), Inches(2.6),
             [("① AI 不是「收到命令→执行」的助手", 17, WHITE, True, True),
              ("而是「理解 → 建模 → 规划 → 执行 → 验证 → 学习」的执行者", 15, MUTED, False, False),
              ("② 按角色而非工具定义 · 文件即数据库 · 治理先于执行", 17, WHITE, True, True),
              ("③ 自动化有护栏：低风险自动，高风险人审，超限升级人工", 17, WHITE, True, True),
              ("", 6, MUTED, False, False),
              ("一个设想：先别问框架能不能做好软件 —— 它已经做好了一个（aimonitor）", 16, CYAN, True, True)],
             spacing=1.2, space_after=10)
    add_badge(s, Inches(0.6), Inches(5.6), Inches(2.6), Inches(0.6), "Q & A", RED, size=18)
    add_text(s, Inches(3.6), Inches(5.7), Inches(8), Inches(0.5),
             "谢谢聆听 · 欢迎提问", size=18, color=MUTED)
    add_text(s, Inches(0.6), Inches(6.5), Inches(12), Inches(0.4),
             "github: 你的仓库地址 · 作者: hb · MIT License", size=12, color=GRID)

# ================= 主流程 =================

def main():
    out = os.environ.get("PPT_OUT", "docs/PPT-AIBASE-INTRO.pptx")
    if len(__import__("sys").argv) > 1:
        out = __import__("sys").argv[1]
    prs = new_prs()
    builders = [slide_01, slide_02, slide_03, slide_04, slide_05, slide_06,
                slide_07, slide_08, slide_09, slide_10, slide_11, slide_12,
                slide_13, slide_14, slide_15, slide_16, slide_17, slide_18,
                slide_19, slide_20]
    for b in builders:
        b(prs)
    os.makedirs(os.path.dirname(os.path.abspath(out)) or ".", exist_ok=True)
    prs.save(out)
    n = len(prs.slides._sldIdLst)
    print(f"✓ 已生成 {out}（{n} 页）")
    assert n == TOTAL, f"页数错误: {n} != {TOTAL}"

if __name__ == "__main__":
    main()
